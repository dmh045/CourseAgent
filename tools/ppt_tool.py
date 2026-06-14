from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List


SLIDE_W = 13.333
SLIDE_H = 7.5


def create_ppt(
    outline: Dict[str, Any],
    scripts: Dict[str, Any],
    mindmap_image_path: str | None,
    output_path: str,
) -> str:
    """Create a polished classroom-defense PPT from outline and scripts."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise ValueError("缺少 python-pptx 依赖，无法生成 PPT。") from exc

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    scripts_by_page = {item.get("page"): item for item in scripts.get("scripts", [])}
    slides = outline.get("slides", [])
    topic = str(outline.get("title") or _topic_from_slides(slides))

    _add_cover(prs, topic, len(slides))

    if mindmap_image_path and Path(mindmap_image_path).suffix.lower() == ".png" and Path(mindmap_image_path).exists():
        _add_mindmap_slide(prs, topic, slides)

    for index, item in enumerate(slides, start=1):
        script = scripts_by_page.get(item.get("page"), {}).get("script", "")
        _add_content_slide(prs, item, script, index, len(slides), topic)

    prs.save(str(output))
    return str(output)


def _add_cover(prs: Any, topic: str, slide_count: int) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "F8FAFC")
    _shape(slide, 0, 0, SLIDE_W, 7.5, "0F172A")
    _shape(slide, 0, 5.65, SLIDE_W, 1.85, "1D4ED8")
    _shape(slide, 0.72, 0.72, 0.12, 4.5, "38BDF8")

    eyebrow = slide.shapes.add_textbox(Inches(1.05), Inches(0.92), Inches(4.4), Inches(0.35))
    _set_text(eyebrow.text_frame, "课程资料整理 / 课堂答辩", 14, "BAE6FD", bold=True)

    title_box = slide.shapes.add_textbox(Inches(1.02), Inches(1.62), Inches(10.9), Inches(1.45))
    frame = title_box.text_frame
    frame.word_wrap = True
    frame.text = topic
    para = frame.paragraphs[0]
    para.font.name = "Microsoft YaHei"
    para.font.bold = True
    para.font.size = Pt(40 if len(topic) <= 16 else 34)
    para.font.color.rgb = RGBColor(248, 250, 252)

    subtitle = slide.shapes.add_textbox(Inches(1.05), Inches(3.25), Inches(9.8), Inches(0.75))
    _set_text(subtitle.text_frame, "从原始资料提炼知识结构、答辩大纲、讲稿和讲解视频", 20, "CBD5E1")

    chips = [
        f"{slide_count} 页内容结构",
        "含思维导图",
        "含语音讲解视频",
    ]
    for idx, text in enumerate(chips):
        x = 1.05 + idx * 2.42
        chip = _shape(slide, x, 4.38, 2.08, 0.42, "E0F2FE", radius=True)
        _set_text(chip.text_frame, text, 12, "075985", bold=True, align=PP_ALIGN.CENTER)

    footer = slide.shapes.add_textbox(Inches(1.05), Inches(6.24), Inches(11.1), Inches(0.5))
    _set_text(footer.text_frame, "自动生成的 PPT 已保留可编辑文本，并将完整讲稿写入演讲者备注。", 14, "EFF6FF", align=PP_ALIGN.CENTER)


def _add_mindmap_slide(prs: Any, topic: str, slides: List[Dict[str, Any]]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "F8FAFC")
    _add_top_chrome(slide, "知识结构总览", "思维导图", topic)
    _shape(slide, 0.72, 1.18, 11.9, 5.76, "FFFFFF", line="DBEAFE", radius=True)

    center = _shape(slide, 5.12, 3.02, 3.1, 1.0, "0F172A", radius=True)
    center.text_frame.word_wrap = True
    _set_text(center.text_frame, topic, 20 if len(topic) <= 12 else 16, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    branch_items = _mindmap_branches(slides)
    positions = [
        (1.04, 1.52),
        (1.04, 3.04),
        (1.04, 4.56),
        (8.56, 1.52),
        (8.56, 3.04),
        (8.56, 4.56),
    ]
    colors = ["1D4ED8", "0284C7", "059669", "7C3AED", "EA580C", "475569"]
    for idx, branch in enumerate(branch_items[:6]):
        x, y = positions[idx]
        accent = colors[idx % len(colors)]
        card = _shape(slide, x, y, 3.42, 1.12, "F8FAFC", line="E2E8F0", radius=True)
        _shape(slide, x, y, 0.12, 1.12, accent)
        title_box = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.12), Inches(3.02), Inches(0.33))
        _set_text(title_box.text_frame, branch["title"], 13, accent, bold=True)
        body_box = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.48), Inches(3.0), Inches(0.48))
        body_box.text_frame.word_wrap = True
        _set_text(body_box.text_frame, branch["detail"], 10, "334155")

        line_start_x = x + 3.42 if x < 5 else x
        line_end_x = 5.12 if x < 5 else 8.22
        line_y = y + 0.56
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(line_start_x),
            Inches(line_y),
            Inches(line_end_x),
            Inches(3.52),
        )
        connector.line.color.rgb = RGBColor.from_string("CBD5E1")
        connector.line.width = Pt(1.2)

    hint = slide.shapes.add_textbox(Inches(4.35), Inches(4.35), Inches(4.0), Inches(0.52))
    _set_text(hint.text_frame, "按这张图讲：先定义模型，再说明符号、时间参数、规则和案例。", 11, "64748B", align=PP_ALIGN.CENTER)
    _add_notes_fallback(slide, "这一页用于建立整体认知。讲解时先说明中心主题，再按左右分支说明资料如何组织知识点，最后引到后续逐页展开。")


def _add_content_slide(prs: Any, item: Dict[str, Any], script: str, index: int, total: int, topic: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "F8FAFC")
    title = str(item.get("title", ""))
    page = item.get("page", index)
    bullets = [str(b).strip() for b in item.get("bullets", []) if str(b).strip()]
    claim = bullets[0] if bullets else title
    details = bullets[1:] if len(bullets) > 1 else bullets

    _add_top_chrome(slide, title, f"第 {page} 页 / {total}", topic)

    claim_box = _shape(slide, 0.82, 1.2, 11.72, 0.92, "E0F2FE", line="7DD3FC", radius=True)
    claim_box.text_frame.word_wrap = True
    _set_text(claim_box.text_frame, claim, 19, "075985", bold=True)

    card_top = 2.34
    card_h = 0.82 if len(details) >= 5 else 0.94
    gap_y = 0.18
    col_w = 5.68
    positions = []
    for idx in range(max(len(details), 1)):
        col = idx % 2
        row = idx // 2
        positions.append((0.82 + col * 6.04, card_top + row * (card_h + gap_y)))

    for idx, detail in enumerate(details[:6], start=1):
        x, y = positions[idx - 1]
        card = _shape(slide, x, y, col_w, card_h, "FFFFFF", line="E2E8F0", radius=True)
        _shape(slide, x + 0.16, y + 0.19, 0.34, 0.34, "1D4ED8", radius=True)
        num = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.17), Inches(0.34), Inches(0.22))
        _set_text(num.text_frame, str(idx), 9, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        body = slide.shapes.add_textbox(Inches(x + 0.62), Inches(y + 0.13), Inches(col_w - 0.78), Inches(card_h - 0.2))
        body.text_frame.word_wrap = True
        _set_text(body.text_frame, detail, _detail_font_size(detail, len(details)), "0F172A")

    if script:
        tip = _speaker_tip(script)
        tip_box = _shape(slide, 0.82, 6.32, 11.72, 0.62, "F1F5F9", line="CBD5E1", radius=True)
        tip_box.text_frame.word_wrap = True
        _set_text(tip_box.text_frame, f"讲解提示：{tip}", 11, "475569")
        _add_notes_fallback(slide, script)

    page_no = slide.shapes.add_textbox(Inches(11.75), Inches(0.43), Inches(0.7), Inches(0.3))
    frame = page_no.text_frame
    frame.text = f"{index:02d}"
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    para.font.name = "Microsoft YaHei"
    para.font.size = Pt(13)
    para.font.bold = True
    para.font.color.rgb = RGBColor(100, 116, 139)


def _add_top_chrome(slide: Any, title: str, label: str, topic: str) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    _shape(slide, 0, 0, 0.26, SLIDE_H, "1D4ED8")
    _shape(slide, 0.26, 0, SLIDE_W - 0.26, 0.18, "38BDF8")
    title_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.38), Inches(8.8), Inches(0.52))
    _set_text(title_box.text_frame, title, 25 if len(title) <= 24 else 21, "0F172A", bold=True)
    tag = _shape(slide, 9.85, 0.44, 1.55, 0.34, "DBEAFE", radius=True)
    _set_text(tag.text_frame, label, 10, "1E40AF", bold=True, align=PP_ALIGN.CENTER)
    topic_box = slide.shapes.add_textbox(Inches(0.82), Inches(6.98), Inches(6.0), Inches(0.24))
    _set_text(topic_box.text_frame, topic, 9, "94A3B8")


def _shape(slide: Any, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, radius: bool = False) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = Inches(0.01)
    else:
        shape.line.fill.background()
    return shape


def _set_bg(slide: Any, color: str) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def _set_text(frame: Any, text: str, size: int, color: str, bold: bool = False, align: Any | None = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    frame.clear()
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(3)
    frame.margin_bottom = Pt(3)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)


def _detail_font_size(text: str, count: int) -> int:
    if len(text) > 58 or count >= 5:
        return 13
    if len(text) > 42:
        return 14
    return 15


def _speaker_tip(script: str) -> str:
    cleaned = " ".join(script.split())
    return cleaned[:120] + ("..." if len(cleaned) > 120 else "")


def _mindmap_branches(slides: List[Dict[str, Any]]) -> List[Dict[str, str]]:
    branches: List[Dict[str, str]] = []
    for item in slides[:6]:
        title = str(item.get("title", "")).strip() or "知识分支"
        short_title = title.split("：", 1)[-1] if "：" in title else title
        bullets = [str(b).strip() for b in item.get("bullets", []) if str(b).strip()]
        detail = "；".join(bullets[:2]) if bullets else title
        branches.append(
            {
                "title": short_title[:20],
                "detail": detail[:54] + ("..." if len(detail) > 54 else ""),
            }
        )
    while len(branches) < 6:
        branches.append({"title": "补充理解", "detail": "结合原文继续拆解概念、规则、案例与适用边界。"})
    return branches


def _add_picture_fit(slide: Any, image_path: str, x: float, y: float, max_w: float, max_h: float) -> None:
    from pptx.util import Inches
    from PIL import Image

    with Image.open(image_path) as image:
        aspect = image.width / max(image.height, 1)
    box_aspect = max_w / max_h
    if aspect >= box_aspect:
        width = max_w
        height = max_w / aspect
    else:
        height = max_h
        width = max_h * aspect
    left = x + (max_w - width) / 2
    top = y + (max_h - height) / 2
    slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def _topic_from_slides(slides: List[Dict[str, Any]]) -> str:
    if not slides:
        return "课程资料主题"
    title = str(slides[0].get("title", "课程资料主题"))
    return title.replace("：主题与背景", "").replace("概览", "").strip(" ：") or title


def _add_notes_fallback(slide: Any, script: str) -> None:
    try:
        notes = slide.notes_slide.notes_text_frame
        notes.text = script
    except Exception:
        return
